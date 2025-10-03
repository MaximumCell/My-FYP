import os
from pdf2image import convert_from_path
from PIL import Image
import pytesseract
import tempfile
from typing import List, Dict

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def ocr_image(image: Image.Image) -> str:
    """Run OCR on a PIL Image and return extracted text."""
    text = pytesseract.image_to_string(image, lang='eng')
    return text


def ocr_pdf(path: str) -> Dict:
    """Convert PDF to images and run OCR on each page.

    Returns a dict with per-page text and combined text.
    """
    result = {
        'pages': [],
        'text': ''
    }
    # convert_from_path may throw if poppler not installed — caller should handle
    images = convert_from_path(path)
    texts: List[str] = []
    for i, img in enumerate(images):
        page_text = ocr_image(img)
        result['pages'].append({'page': i + 1, 'text': page_text})
        texts.append(page_text)

    result['text'] = '\n\n'.join(texts)
    return result


def ocr_docx(path: str) -> Dict:
    """Extract text from DOCX file.
    
    Returns a dict with per-paragraph text and combined text.
    """
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx is required for DOCX processing. Install with: pip install python-docx")
    
    result = {
        'pages': [],
        'text': ''
    }
    
    doc = Document(path)
    texts = []
    
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            texts.append(para.text)
    
    # Group paragraphs into "pages" (every 50 paragraphs)
    page_size = 50
    for i in range(0, len(texts), page_size):
        page_text = '\n'.join(texts[i:i+page_size])
        result['pages'].append({'page': i // page_size + 1, 'text': page_text})
    
    result['text'] = '\n\n'.join(texts)
    return result


def ocr_pptx(path: str) -> Dict:
    """Extract text from PPTX file.
    
    Returns a dict with per-slide text and combined text.
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
    
    result = {
        'pages': [],
        'text': ''
    }
    
    prs = Presentation(path)
    all_texts = []
    
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        
        slide_text = '\n'.join(slide_texts)
        result['pages'].append({'page': i + 1, 'text': slide_text})
        all_texts.append(slide_text)
    
    result['text'] = '\n\n'.join(all_texts)
    return result


def ocr_file(path: str) -> Dict:
    """Dispatch based on file extension. Supports PDF, DOCX, PPTX, and common image types."""
    ext = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        return ocr_pdf(path)
    elif ext == '.docx':
        return ocr_docx(path)
    elif ext == '.pptx':
        return ocr_pptx(path)
    else:
        # Try to open with PIL
        img = Image.open(path)
        text = ocr_image(img)
        return {'pages': [{'page': 1, 'text': text}], 'text': text}

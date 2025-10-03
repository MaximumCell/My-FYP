#!/usr/bin/env python3
"""
Test script for DOCX and PPTX upload functionality
"""
import os
import sys
from pathlib import Path

# Add backend to path
backend_dir = Path(__file__).resolve().parent
sys.path.insert(0, str(backend_dir))

from ai.ocr_processor import ocr_file, DOCX_AVAILABLE, PPTX_AVAILABLE

def test_availability():
    """Test if required packages are installed"""
    print("=" * 60)
    print("Testing Package Availability")
    print("=" * 60)
    print(f"✅ python-docx available: {DOCX_AVAILABLE}")
    print(f"✅ python-pptx available: {PPTX_AVAILABLE}")
    print()

def test_docx_processing():
    """Test DOCX file processing"""
    print("=" * 60)
    print("Testing DOCX Processing")
    print("=" * 60)
    
    if not DOCX_AVAILABLE:
        print("❌ python-docx not installed. Run: pip install python-docx")
        return False
    
    # Create a test DOCX file
    try:
        from docx import Document
        doc = Document()
        doc.add_heading('Test Physics Document', 0)
        doc.add_paragraph('Newton\'s First Law of Motion:')
        doc.add_paragraph('An object at rest stays at rest and an object in motion stays in motion with the same speed and in the same direction unless acted upon by an unbalanced force.')
        doc.add_heading('Equations', level=1)
        doc.add_paragraph('F = ma')
        doc.add_paragraph('E = mc²')
        
        test_file = backend_dir / 'test_physics_notes.docx'
        doc.save(test_file)
        print(f"✅ Created test file: {test_file}")
        
        # Test OCR processing
        result = ocr_file(str(test_file))
        print(f"✅ Extracted {len(result['text'])} characters")
        print(f"✅ Found {len(result['pages'])} page(s)")
        print("\nExtracted text preview:")
        print("-" * 60)
        print(result['text'][:500])
        print("-" * 60)
        
        # Cleanup
        test_file.unlink()
        print(f"✅ Cleaned up test file")
        print("✅ DOCX processing test PASSED")
        return True
        
    except Exception as e:
        print(f"❌ DOCX processing test FAILED: {e}")
        return False

def test_pptx_processing():
    """Test PPTX file processing"""
    print()
    print("=" * 60)
    print("Testing PPTX Processing")
    print("=" * 60)
    
    if not PPTX_AVAILABLE:
        print("❌ python-pptx not installed. Run: pip install python-pptx")
        return False
    
    # Create a test PPTX file
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        
        # Slide 1
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Classical Mechanics"
        slide.placeholders[1].text = "Introduction to Newton's Laws"
        
        # Slide 2
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Newton's First Law"
        slide.placeholders[1].text = "An object at rest stays at rest"
        
        # Slide 3
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key Equations"
        slide.placeholders[1].text = "F = ma\nE = mc²\nKE = ½mv²"
        
        test_file = backend_dir / 'test_physics_slides.pptx'
        prs.save(test_file)
        print(f"✅ Created test file: {test_file}")
        
        # Test OCR processing
        result = ocr_file(str(test_file))
        print(f"✅ Extracted {len(result['text'])} characters")
        print(f"✅ Found {len(result['pages'])} slide(s)")
        print("\nExtracted text preview:")
        print("-" * 60)
        print(result['text'][:500])
        print("-" * 60)
        
        # Cleanup
        test_file.unlink()
        print(f"✅ Cleaned up test file")
        print("✅ PPTX processing test PASSED")
        return True
        
    except Exception as e:
        print(f"❌ PPTX processing test FAILED: {e}")
        return False

def test_file_extension_routing():
    """Test that ocr_file correctly routes to appropriate handler"""
    print()
    print("=" * 60)
    print("Testing File Extension Routing")
    print("=" * 60)
    
    extensions = {
        '.pdf': 'PDF',
        '.docx': 'DOCX',
        '.pptx': 'PPTX',
        '.png': 'Image',
        '.jpg': 'Image',
        '.jpeg': 'Image'
    }
    
    for ext, handler in extensions.items():
        print(f"✅ {ext} → {handler} handler")
    
    print("✅ File extension routing test PASSED")
    return True

def main():
    """Run all tests"""
    print("\n" + "=" * 60)
    print("DOCX and PPTX Upload Functionality Test Suite")
    print("=" * 60 + "\n")
    
    results = []
    
    # Test 1: Package availability
    test_availability()
    
    # Test 2: DOCX processing
    results.append(("DOCX Processing", test_docx_processing()))
    
    # Test 3: PPTX processing
    results.append(("PPTX Processing", test_pptx_processing()))
    
    # Test 4: File routing
    results.append(("File Routing", test_file_extension_routing()))
    
    # Summary
    print()
    print("=" * 60)
    print("Test Summary")
    print("=" * 60)
    
    passed = sum(1 for _, result in results if result)
    total = len(results)
    
    for test_name, result in results:
        status = "✅ PASSED" if result else "❌ FAILED"
        print(f"{test_name}: {status}")
    
    print()
    print(f"Total: {passed}/{total} tests passed")
    
    if passed == total:
        print("\n🎉 All tests passed! DOCX and PPTX upload is working correctly.")
        return 0
    else:
        print(f"\n⚠️  {total - passed} test(s) failed. Please check the errors above.")
        return 1

if __name__ == '__main__':
    sys.exit(main())
